import json
import os
import re
from pptx import Presentation
from PyPDF2 import PdfReader
import docx

def clean_text(text):
    if not text:
        return "missing_from_source"
    cleaned = re.sub(r'\s+', ' ', text).strip()
    return cleaned if cleaned else "missing_from_source"

def parse_docx(filepath):
    """Returns a dict mapping module_id to simulation data."""
    if not os.path.exists(filepath):
        return {}
    
    doc = docx.Document(filepath)
    sims = {}
    current_module = None
    lines = [p.text for p in doc.paragraphs if p.text.strip()]
    
    for i, line in enumerate(lines):
        # Look for "Module X.Y"
        mod_match = re.search(r'Module\s+(\d+\.\d+)', line, re.IGNORECASE)
        if mod_match:
            current_module = mod_match.group(1)
            sims[current_module] = {
                "simulation_scenario_name": "missing_from_source",
                "prospect_persona": "missing_from_source",
                "difficulty_level": "missing_from_source",
                "sales_skill_practiced": "missing_from_source",
                "simulation_goal": "missing_from_source",
                "_raw": ""
            }
        
        if current_module:
            sims[current_module]["_raw"] += line + "\n"

    for mod, data in sims.items():
        raw = data["_raw"]
        
        persona_match = re.search(r'(?:AI|Customer|Prospect) Persona:\s*(.*?)\n', raw, re.IGNORECASE)
        if persona_match: data["prospect_persona"] = persona_match.group(1).strip()
        
        diff_match = re.search(r'Difficulty(?:\s*Level)?:\s*(.*?)\n', raw, re.IGNORECASE)
        if diff_match: data["difficulty_level"] = diff_match.group(1).strip()
        
        skills_match = re.search(r'(?:Primary )?Skill[s]?(?: Trained| Practiced)?:\s*(.*?)\n', raw, re.IGNORECASE)
        if skills_match: data["sales_skill_practiced"] = skills_match.group(1).strip()
        
        goal_match = re.search(r'(?:Scenario )?Goal:\s*(.*?)\n', raw, re.IGNORECASE)
        if goal_match: data["simulation_goal"] = goal_match.group(1).strip()
        
        name_match = re.search(r'(?:Scenario )?(?:Name|Title):\s*(.*?)\n', raw, re.IGNORECASE)
        if name_match: 
            data["simulation_scenario_name"] = name_match.group(1).strip()
        else:
            data["simulation_scenario_name"] = f"{mod} Exercise"

    return sims

def parse_pptx(filepath):
    if not os.path.exists(filepath):
        return []
    
    prs = Presentation(filepath)
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        title = "missing_from_source"
        if slide.shapes.title and slide.shapes.title.text:
            title = clean_text(slide.shapes.title.text)
        
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes.title:
                if shape.text:
                    text_content.append(shape.text)
        
        content = clean_text(" ".join(text_content))
        
        notes = "missing_from_source"
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = clean_text(slide.notes_slide.notes_text_frame.text)
            
        slides_data.append({
            "slide_number": i + 1,
            "title": title,
            "content": content,
            "notes": notes
        })
        
    return slides_data

def parse_pdf(filepath):
    if not os.path.exists(filepath):
        return []
    
    try:
        reader = PdfReader(filepath)
        slides_data = []
        for i, page in enumerate(reader.pages):
            text = clean_text(page.extract_text())
            # Crude separation of title vs content assuming first few lines are title
            lines = text.split('\n')
            title = lines[0] if lines else "missing_from_source"
            content = " ".join(lines[1:]) if len(lines) > 1 else "missing_from_source"
            
            slides_data.append({
                "slide_number": i + 1,
                "title": title[:100], # Keep title short
                "content": content,
                "notes": "missing_from_source"
            })
        return slides_data
    except Exception as e:
        print(f"Error reading PDF {filepath}: {e}")
        return []

def main():
    json_path = os.path.join(os.path.dirname(__file__), "curriculum_map.json")
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
        
    base_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
    
    for day_obj in data["days"]:
        day_num = day_obj["day_number"]
        
        # 1. Update Simulations from DOCX
        docx_path = os.path.join(base_dir, "AntiGravity Doc", "training_materials", "files (1)", f"Day{day_num}_Solar_Sales_Accelerator_AI_Simulation.docx")
        sims = parse_docx(docx_path)
        
        # 2. Update Slides from PPTX/PDF
        pptx_path = os.path.join(base_dir, "solar-trainer", f"SolarAccelerator_Day{day_num}.pptx")
        pdf_path = pptx_path + ".pdf"
        
        if os.path.exists(pptx_path):
            slides_extracted = parse_pptx(pptx_path)
        elif os.path.exists(pdf_path):
            slides_extracted = parse_pdf(pdf_path)
        else:
            slides_extracted = []
            
        # Distribute slides to modules based on slide title containing "Module X.Y"
        # If no explicit module in title, we continue adding to the current module.
        current_module_idx = 0
        
        if slides_extracted and day_obj["modules"]:
            # Clear existing placeholder slides for all modules
            for mod in day_obj["modules"]:
                mod["slides"] = []
                mod["transcript_segments"] = []
                mod["number_of_slides"] = 0
                mod["trainer_script_sections"] = 0
                
            for slide in slides_extracted:
                title = slide["title"]
                
                # Check if this slide signals a new module
                for m_idx, mod in enumerate(day_obj["modules"]):
                    if mod["module_id"] in title or f"Module {mod['module_id']}" in title:
                        current_module_idx = m_idx
                        break
                        
                target_mod = day_obj["modules"][current_module_idx]
                
                # Append slide
                new_slide_num = len(target_mod["slides"]) + 1
                target_mod["slides"].append({
                    "slide_number": new_slide_num,
                    "slide_title": title,
                    "slide_summary": slide["content"][:200] if slide["content"] != "missing_from_source" else "missing_from_source",
                    "key_concept": slide["title"],
                    "supporting_trainer_script_segment": slide["notes"]
                })
                
                # Append transcript
                if slide["notes"] != "missing_from_source":
                    target_mod["transcript_segments"].append({
                        "module_id": target_mod["module_id"],
                        "slide_number": new_slide_num,
                        "transcript_segment": slide["notes"],
                        "key_teaching_point": slide["title"]
                    })
                    target_mod["transcript_ready"] = True
                
            # Update counts
            for mod in day_obj["modules"]:
                mod["number_of_slides"] = len(mod["slides"])
                mod["trainer_script_sections"] = len(mod["transcript_segments"])

        # Distribute simulations
        for mod in day_obj["modules"]:
            mod_id = mod["module_id"]
            if mod_id in sims:
                sim_data = sims[mod_id]
                mod["simulation_mapping"]["simulation_scenario_name"] = sim_data.get("simulation_scenario_name", "missing_from_source")
                mod["simulation_mapping"]["prospect_persona"] = sim_data.get("prospect_persona", "missing_from_source")
                mod["simulation_mapping"]["difficulty_level"] = sim_data.get("difficulty_level", "missing_from_source")
                mod["simulation_mapping"]["sales_skill_practiced"] = sim_data.get("sales_skill_practiced", "missing_from_source")
                mod["simulation_mapping"]["simulation_goal"] = sim_data.get("simulation_goal", "missing_from_source")
                mod["simulation_ready"] = True

    with open(json_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=2)
        
    print("Done updating curriculum map.")

if __name__ == "__main__":
    main()
