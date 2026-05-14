"""
inject_images.py
Inserts AI-generated images into Solar_Sales_Training_Day2.pptx
based on the Day2_ImagePromptGuide placement coordinates.
Slide width=10.0", height=5.625" (16:9 widescreen)
"""
import os
import io
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

ROOT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
NEW_DIR  = os.path.join(ROOT_DIR, "docs", "AntiGravity Doc", "training_materials")
PPTX_IN  = os.path.join(NEW_DIR, "Solar_Sales_Training_Day2.pptx")
PPTX_OUT = os.path.join(NEW_DIR, "Solar_Sales_Training_Day2_Final.pptx")

# (slide_index_0based, filename, left_in, top_in, width_in, height_in, send_to_back, crop_strip)
# crop_strip=True  -> center-crop image to match the slot's aspect ratio before inserting
PLACEMENTS = [
    # Slide 1  – Title hero: right half background
    (0,  "solar_tech_roof_hero.png.png",       5.0,  0.0,   5.0,  5.625, True,  False),
    # Slide 3  – Module 2.1 Divider: solar field, right half background
    (2,  "solar_field_golden_hour.png.png",    5.0,  0.0,   5.0,  5.625, True,  False),
    # Slide 5  – Reality Check: homeowner at door, right column panel
    (4,  "homeowner_door_guarded.png.png",     5.2,  2.0,   4.5,  2.5,   False, False),
    # Slide 7  – Decision Avoidance: couch, bottom strip (auto-cropped)
    (6,  "status_quo_couch.png (2).png",       0.1,  4.47,  9.79, 1.15,  False, True),
    # Slide 8  – Contractor Trust Barrier: arms crossed, right column
    (7,  "homeowner_arms_crossed.png.png",     6.1,  0.9,   3.6,  3.6,   True,  False),
    # Slide 9  – Addressing Skepticism: handshake, banner background (auto-cropped)
    (8,  "handshake_porch.png (2).png",        0.1,  0.85,  9.79, 1.25,  True,  True),
    # Slide 12 – Guarded Gloria: left scenario panel
    (11, "guarded_gloria.png.png",             0.3,  0.85,  4.5,  1.7,   False, False),
    # Slide 13 – Module 2.2 Divider: aerial suburb, right half background
    (12, "aerial_suburb.png.png",              5.0,  0.0,   5.0,  5.625, True,  False),
    # Slide 16 – Prospecting Tools: rep at door, bottom strip (auto-cropped)
    (15, "solar_rep_door_back.png.png",        0.1,  4.79,  9.79, 0.83,  False, True),
    # Slide 27 – 30-Second Pitch: rep at door (full right column)
    (26, "solar_rep_door_pitch.png.png",       5.2,  0.88,  4.5,  4.42,  False, False),
    # Slide 36 – Anti-Pitch Strategy: consultant at table, banner inset (auto-cropped)
    (35, "consultant_kitchen_table.png.png",   6.4,  0.85,  3.3,  1.0,   False, True),
    # Slide 44 – Module 2.11 Divider: rep walking, right half background
    (43, "solar_rep_walking.png.png",          5.0,  0.0,   5.0,  5.625, True,  False),
]


def center_crop(img_path, target_w_in, target_h_in):
    """Center-crop a PIL image to match the target aspect ratio. Returns a BytesIO PNG."""
    img = Image.open(img_path)
    src_w, src_h = img.size
    target_ar = target_w_in / target_h_in

    # Determine crop box
    if src_w / src_h > target_ar:
        # Image is wider than needed — crop width
        new_w = int(src_h * target_ar)
        x0 = (src_w - new_w) // 2
        box = (x0, 0, x0 + new_w, src_h)
    else:
        # Image is taller than needed — crop height
        new_h = int(src_w / target_ar)
        y0 = (src_h - new_h) // 2
        box = (0, y0, src_w, y0 + new_h)

    cropped = img.crop(box)
    buf = io.BytesIO()
    cropped.save(buf, format="PNG")
    buf.seek(0)
    return buf


def send_to_back(slide, shape):
    """Move shape to behind all other shapes on the slide."""
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def main():
    prs = Presentation(PPTX_IN)
    print(f"Loaded: {len(prs.slides)} slides  ({prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\")\n")

    inserted = 0
    skipped  = 0

    for slide_idx, img_name, left, top, width, height, to_back, crop_strip in PLACEMENTS:
        img_path = os.path.join(NEW_DIR, img_name)
        if not os.path.exists(img_path):
            print(f"  SKIP  slide {slide_idx+1:2d} – not found: {img_name}")
            skipped += 1
            continue

        slide = prs.slides[slide_idx]

        if crop_strip:
            img_data = center_crop(img_path, width, height)
            pic = slide.shapes.add_picture(img_data, Inches(left), Inches(top), Inches(width), Inches(height))
            crop_note = "  [strip-cropped]"
        else:
            pic = slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
            crop_note = ""

        if to_back:
            send_to_back(slide, pic)
            print(f"  OK    slide {slide_idx+1:2d} – {img_name}  [behind text]{crop_note}")
        else:
            print(f"  OK    slide {slide_idx+1:2d} – {img_name}{crop_note}")
        inserted += 1

    prs.save(PPTX_OUT)
    print(f"\n{inserted} image(s) inserted, {skipped} skipped.")
    print(f"Saved: {PPTX_OUT}")


if __name__ == "__main__":
    main()
